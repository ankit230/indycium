from pptx import Presentation
from pptx.util import Inches
pr1 = Presentation()

slide1_register = pr1.slide_layouts[0]
slide2_register = pr1.slide_layouts[1]
slide3_register = pr1.slide_layouts[1]
slide4_register = pr1.slide_layouts[1]
slide5_register = pr1.slide_layouts[1]
slide6_register = pr1.slide_layouts[1]

slide1 = pr1.slides.add_slide(slide1_register)
slide2 = pr1.slides.add_slide(slide2_register)
slide3 = pr1.slides.add_slide(slide3_register)
slide4 = pr1.slides.add_slide(slide4_register)
slide5 = pr1.slides.add_slide(slide5_register)
slide6 = pr1.slides.add_slide(slide6_register)


title1 = slide1.shapes.title
title2 = slide2.shapes.title
title3 = slide3.shapes.title
title4 = slide4.shapes.title
title5 = slide5.shapes.title
title6 = slide6.shapes.title

subtitle1 = slide1.placeholders[1]
subtitle2 = slide2.placeholders[1]
subtitle3 = slide3.placeholders[1]
subtitle4 = slide4.placeholders[1]
subtitle5 = slide5.placeholders[1]
subtitle6 = slide6.placeholders[1]

title1.text = "INDYCIUM - Assignment"
title2.text = "First Image"
title3.text = "Second Image"
title4.text = "Third Image"
title5.text = "Fourth Image"
title6.text = "Fifth Image"



subtitle1.text = "Powerpoint Presentation Using Python"
subtitle2.text = "Food"
subtitle3.text = "Interior"
subtitle4.text = "Coding"
subtitle5.text = "Photgraphy"
subtitle6.text = "Decor"

img1 = "111.jpg"
img2 = "22.jpg"
img3 = "33.jpg"
img4 = "44.jpg"
img5 = "55.jpg"

from_left = Inches(4)
from_left4 = Inches(3)
from_top2 = Inches(3)

add_picture = slide2.shapes.add_picture(img1,from_left,from_top2)
add_picture = slide3.shapes.add_picture(img2,from_left,from_top2)
add_picture = slide4.shapes.add_picture(img3,from_left4,from_top2)
add_picture = slide5.shapes.add_picture(img4,from_left4,from_top2)
add_picture = slide6.shapes.add_picture(img5,from_left,from_top2)


pr1.save('assignppt.pptx')